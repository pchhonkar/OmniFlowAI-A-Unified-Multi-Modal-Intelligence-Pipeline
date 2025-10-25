import os
import json
import re
import pandas as pd
from pathlib import Path
from dotenv import load_dotenv
from PIL import Image
from docx import Document as WordDoc
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
import matplotlib.pyplot as plt
import requests
from io import BytesIO
import pyttsx3
from moviepy.editor import TextClip, concatenate_videoclips

from duckduckgo_search import DDGS
from langchain.agents import Tool, initialize_agent, AgentType
from langchain.vectorstores import FAISS
from langchain.docstore.document import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_community.utilities.wikipedia import WikipediaAPIWrapper
from langchain_community.document_loaders import (
    PyPDFLoader, 
    TextLoader, 
    Docx2txtLoader, 
    CSVLoader,
    UnstructuredExcelLoader, 
    UnstructuredMarkdownLoader
)

# === Load env ===
load_dotenv()
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY", "")
os.environ["HF_TOKEN"] = os.getenv("HF_TOKEN", "")
os.environ["HUGGINGFACEHUB_API_TOKEN"] = os.getenv("HUGGINGFACEHUB_API_TOKEN", os.getenv("HF_TOKEN", ""))

# === Folder Setup ===
for folder in ["data", "Output"]:
    Path(folder).mkdir(parents=True, exist_ok=True)

# === Extract text from PPTX ===
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# === Load all documents ===
def load_documents(folder="data"):
    docs = []
    for file in os.listdir(folder):
        path = os.path.join(folder, file)
        ext = os.path.splitext(file)[1].lower()
        try:
            if ext == ".pdf":
                docs.extend(PyPDFLoader(path).load())
            elif ext == ".txt":
                docs.extend(TextLoader(path).load())
            elif ext == ".docx":
                docs.extend(Docx2txtLoader(path).load())
            elif ext == ".csv":
                docs.extend(CSVLoader(path).load())
            elif ext in [".xls", ".xlsx"]:
                docs.extend(UnstructuredExcelLoader(path).load())
            elif ext == ".md":
                docs.extend(UnstructuredMarkdownLoader(path).load())
            elif ext == ".pptx":
                text = extract_text_from_pptx(path)
                docs.append(Document(page_content=text, metadata={"source": file}))
        except Exception as e:
            print(f"❌ Failed to load {file}: {e}")
    return docs

# === FAISS Vectorstore Setup ===
def setup_vectorstore():
    docs = load_documents()
    chunks = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100).split_documents(docs)
    return FAISS.from_documents(chunks, OpenAIEmbeddings())

# === Document and External Search Tools ===
def document_search(query):
    vectorstore_path = Path("vectorstore")
    if not vectorstore_path.exists():
        vs = setup_vectorstore()
        vs.save_local(str(vectorstore_path))
    vs = FAISS.load_local(str(vectorstore_path), OpenAIEmbeddings(), allow_dangerous_deserialization=True)
    return "\n\n".join([d.page_content for d in vs.as_retriever().get_relevant_documents(query)])

def stats_lookup(query):
    output = []
    for file in os.listdir("data"):
        path = os.path.join("data", file)
        if file.endswith(".xlsx"):
            df = pd.read_excel(path)
            matches = df[df.apply(lambda row: query.lower() in row.astype(str).str.lower().to_string(), axis=1)]
            if not matches.empty:
                output.append(matches.to_string())
        elif file.endswith(".json"):
            with open(path, "r") as f:
                for line in f:
                    if query.lower() in line.lower():
                        output.append(line.strip())
    return "\n---\n".join(output) or "No matching stats found."

def external_search(query):
    results = []
    with DDGS() as ddgs:
        for r in ddgs.text(query, max_results=3):
            results.append(f"{r['title']}: {r['href']}")
    return "\n".join(results) if results else "No relevant results found."

def wiki_search(query):
    return WikipediaAPIWrapper().run(query)

def pick_template(request):
    r = request.lower()
    if "grant" in r:
        return "Use grant format: intro, need, goals, impact, budget."
    elif "blog" in r:
        return "6 paragraphs: hook, story, stats, quote, mission, close."
    elif "slide" in r or "presentation" in r:
        return "Make 5-slide PPT with: Title, Bullets, Chart, Image, Summary."
    elif "social" in r or "instagram" in r:
        return "Make 3 Instagram captions with emojis, one blog carousel, one image quote."
    return "Use general format: summary, data, conclusion."

# === Chart/Image Generation ===
def generate_chart_image():
    fig, ax = plt.subplots()
    ax.plot([1, 2, 3], [3, 2, 5], marker='o')
    ax.set_title("Sample Chart")
    ax.set_xlabel("X-axis")
    ax.set_ylabel("Y-axis")
    img_path = "Output/chart.png"
    plt.savefig(img_path)
    plt.close()
    return img_path

def get_image_from_web(query="food bank"):
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=1))
            img_url = results[0]["image"]
            response = requests.get(img_url)
            img = Image.open(BytesIO(response.content))
            img_path = "Output/slide_image.png"
            img.save(img_path)
            return img_path
    except Exception as e:
        print(f"❌ Failed to fetch image: {e}")
        return None

# === Save Output ===
def save_output(text, name="CAFBrain_Agent_Output", outdir="Output", slide_count=5, audio_secs=None, video_secs=None):
    os.makedirs(outdir, exist_ok=True)

    with open(f"{outdir}/{name}.txt", "w", encoding="utf-8") as f:
        f.write(text)

    doc = WordDoc()
    doc.add_heading("CAFBrain Agent Output", 0)
    doc.add_paragraph(text)
    doc.save(f"{outdir}/{name}.docx")

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for line in text.split("\n"):
        try:
            pdf.multi_cell(0, 10, line.encode("latin-1", "replace").decode("latin-1"))
        except:
            continue
    pdf.output(f"{outdir}/{name}.pdf")

    with open(f"{outdir}/{name}.json", "w", encoding="utf-8") as f:
        json.dump({"output": text}, f, indent=2)

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = name.replace("_", " ")

    slides = []
    current = []
    for line in text.split("\n"):
        line = line.strip()
        if re.match(r"^(Slide|slide) \d+[-:]|^•?\s*\d+\.\s*Slide", line):
            if current:
                slides.append(current)
            current = [line]
        elif line:
            current.append(line)
    if current:
        slides.append(current)
    slides = slides[:slide_count]

    for i, slide_lines in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_lines[0]
        slide.placeholders[1].text = "\n".join(f"• {l}" for l in slide_lines[1:])
    prs.save(f"{outdir}/{name}.pptx")
    print(f"✅ PowerPoint saved at: {outdir}/{name}.pptx")

    if audio_secs:
        try:
            tts = pyttsx3.init()
            tts.setProperty("rate", 160)
            audio_path = f"{outdir}/{name}.mp3"
            tts.save_to_file(text[:4000], audio_path)
            tts.runAndWait()
            print(f"✅ Audio saved at: {audio_path}")
        except Exception as e:
            print(f"❌ TTS error: {e}")

    if video_secs:
        try:
            duration = video_secs // max(1, len(slides))
            clips = []
            for i, chunk in enumerate(slides):
                txt = " | ".join(chunk)
                clip = TextClip(txt, fontsize=24, color="white", size=(1280, 720), method='caption')
                clip = clip.set_duration(duration)
                clips.append(clip)
            final_clip = concatenate_videoclips(clips, method="compose")
            video_path = f"{outdir}/{name}.mp4"
            final_clip.write_videofile(video_path, fps=24)
            print(f"✅ Video saved at: {video_path}")
        except Exception as e:
            print(f"❌ Video error: {e}")

# === Initialize Agent ===
def initialize_cafbrain_agent():
    vectorstore_path = Path("vectorstore").resolve()
    if not vectorstore_path.exists():
        print("🧠 Vectorstore not found. Creating FAISS index...")
        vs = setup_vectorstore()
        vs.save_local(str(vectorstore_path))

    vectorstore = FAISS.load_local(
        str(vectorstore_path),
        OpenAIEmbeddings(),
        allow_dangerous_deserialization=True
    )
    retriever = vectorstore.as_retriever()

    tools = [
        Tool(name="DocumentSearch", func=document_search,
             description="Search CAFB internal documents for facts, statistics, or plans."),
        Tool(name="StatLookup", func=stats_lookup,
             description="Use this to extract budget or metrics from spreadsheets or JSON."),
        Tool(name="TemplatePicker", func=pick_template,
             description="Picks a layout template for grant/blog/slide output."),
        Tool(name="WikiSearch", func=wiki_search,
             description="Get background knowledge from Wikipedia."),
        Tool(name="ExternalSearch", func=external_search,
             description="Pull relevant live web results using DuckDuckGo.")
    ]

    agent_prompt = (
        "You are CAFBrain Agent, an expert grant-writing assistant for Capital Area Food Bank.\n"
        "You are tuned with structured output examples from real grant proposals, blogs, and decks.\n\n"
        "Instructions:\n"
        "1. Think step-by-step using reasoning (Chain-of-Thought).\n"
        "2. Retrieve useful statistics from spreadsheets or PDFs.\n"
        "3. Recommend output structure using an appropriate template.\n"
        "4. Always produce clear, professional text aligned with nonprofit tone.\n\n"
        "Let’s begin.\n"
    )

    llm = ChatOpenAI(model="gpt-4", temperature=0.3)

    agent = initialize_agent(
        tools=tools,
        llm=llm,
        agent=AgentType.ZERO_SHOT_REACT_DESCRIPTION,
        verbose=True,
        handle_parsing_errors=True,
        agent_kwargs={"prefix": agent_prompt}
    )

    return agent

